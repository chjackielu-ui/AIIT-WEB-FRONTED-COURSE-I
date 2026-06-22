#!/usr/bin/env python3
"""生成 舌尖BOSS 项目汇报 PPTX"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ========== 颜色常量 ==========
RED = RGBColor(0xC4, 0x1E, 0x24)
RED_DARK = RGBColor(0x9E, 0x18, 0x20)
RED_LIGHT = RGBColor(0xE8, 0x43, 0x4A)
GOLD = RGBColor(0xD4, 0xA8, 0x43)
GOLD_LIGHT = RGBColor(0xE8, 0xC8, 0x70)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_LIGHT = RGBColor(0x2D, 0x2D, 0x44)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_WARM = RGBColor(0xFF, 0xF8, 0xF0)
BG_LIGHT = RGBColor(0xFF, 0xF0, 0xE0)
CODE_BG = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT = RGBColor(0xCD, 0xD6, 0xF4)
CODE_KEYWORD = RGBColor(0xCB, 0xA6, 0xF7)
CODE_STRING = RGBColor(0xA6, 0xE3, 0xA1)
CODE_COMMENT = RGBColor(0x6C, 0x70, 0x86)
CODE_FUNC = RGBColor(0x89, 0xB4, 0xFA)
GRAY = RGBColor(0x66, 0x66, 0x66)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)

# ========== 创建演示文稿 ==========
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ========== 辅助函数 ==========
def add_bg(slide, color=BG_WARM):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    """添加矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_rounded_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """添加圆角矩形"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_code_block(slide, left, top, width, height, code_text):
    """添加代码块（深色背景+等宽字体）"""
    # 背景矩形
    bg_shape = add_rounded_rect(slide, left, top, width, height, fill_color=CODE_BG)
    bg_shape.shadow.inherit = False
    
    # 代码文本
    txBox = slide.shapes.add_textbox(
        left + Inches(0.3), top + Inches(0.2),
        width - Inches(0.6), height - Inches(0.4)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    
    lines = code_text.strip().split('\n')
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        # 简单语法高亮
        stripped = line.strip()
        if stripped.startswith('//'):
            p.text = line
            p.font.color.rgb = CODE_COMMENT
        elif any(stripped.startswith(kw) for kw in ['const ', 'let ', 'var ', 'function ', 'return ', 'if ', 'else']):
            p.text = line
            p.font.color.rgb = CODE_KEYWORD
        elif stripped.startswith("'") or stripped.startswith('"'):
            p.text = line
            p.font.color.rgb = CODE_STRING
        else:
            p.text = line
            p.font.color.rgb = CODE_TEXT
        
        p.font.size = Pt(13)
        p.font.name = 'Consolas'
        p.space_after = Pt(2)
        p.space_before = Pt(0)
    
    return bg_shape, txBox

def add_decorated_title(slide, title_text, subtitle_text=""):
    """添加带红色左边框装饰的标题"""
    # 左侧红色竖条
    add_rect(slide, Inches(0.6), Inches(0.5), Inches(0.08), Inches(0.65), fill_color=RED)
    # 标题文字
    add_text(slide, Inches(0.85), Inches(0.45), Inches(8), Inches(0.5),
             title_text, font_size=28, color=DARK, bold=True)
    if subtitle_text:
        add_text(slide, Inches(0.85), Inches(0.95), Inches(8), Inches(0.35),
                 subtitle_text, font_size=14, color=GRAY)

def add_page_number(slide, num, total=10):
    """添加页码"""
    add_text(slide, Inches(12.2), Inches(7.0), Inches(0.8), Inches(0.35),
             f"{num}/{total}", font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.RIGHT)

def add_top_bar(slide):
    """添加顶部金色细线装饰"""
    add_rect(slide, Inches(0), Inches(0), W, Inches(0.04), fill_color=GOLD)

def add_bottom_bar(slide):
    """添加底部红色细线装饰"""
    add_rect(slide, Inches(0.5), Inches(6.95), Inches(12.333), Inches(0.025), fill_color=RED)

def add_tag(slide, left, top, text, fill_color=RED, text_color=WHITE):
    """添加标签"""
    w = Inches(0.15 * len(text) + 0.3)
    shape = add_rounded_rect(slide, left, top, w, Inches(0.35), fill_color=fill_color)
    shape.text_frame.paragraphs[0].text = text
    shape.text_frame.paragraphs[0].font.size = Pt(11)
    shape.text_frame.paragraphs[0].font.color.rgb = text_color
    shape.text_frame.paragraphs[0].font.bold = True
    shape.text_frame.paragraphs[0].font.name = 'Microsoft YaHei'
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    return shape

def add_feature_card(slide, left, top, width, height, icon, title, desc, fill_color=WHITE, line_color=GOLD):
    """添加特性卡片"""
    card = add_rounded_rect(slide, left, top, width, height, fill_color=fill_color, line_color=line_color)
    card.line.width = Pt(1.5)
    # 图标
    add_text(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.5), Inches(0.4),
             icon, font_size=22, alignment=PP_ALIGN.CENTER)
    # 标题
    add_text(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), Inches(0.35),
             title, font_size=14, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text(slide, left + Inches(0.15), top + Inches(0.9), width - Inches(0.3), height - Inches(1.1),
             desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    return card

def add_bullet_points(slide, left, top, width, height, items, font_size=14, color=DARK, spacing=Pt(8)):
    """添加要点列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Microsoft YaHei'
        p.space_after = spacing
        p.level = 0
    return txBox

# ============================================================
# 第1页：封面
# ============================================================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
add_bg(slide1, DARK)

# 装饰：大红色背景块
add_rect(slide1, Inches(0), Inches(0), W, Inches(0.12), fill_color=RED)
add_rect(slide1, Inches(0), Inches(7.38), W, Inches(0.12), fill_color=RED)

# 金色装饰线
add_rect(slide1, Inches(3.5), Inches(2.0), Inches(6.333), Inches(0.03), fill_color=GOLD)
add_rect(slide1, Inches(3.5), Inches(5.6), Inches(6.333), Inches(0.03), fill_color=GOLD)

# 项目标题
add_text(slide1, Inches(0), Inches(2.3), W, Inches(1.2),
         "舌尖BOSS", font_size=60, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 副标题
add_text(slide1, Inches(0), Inches(3.5), W, Inches(0.6),
         "中国美食文化探索网站", font_size=26, color=GOLD, alignment=PP_ALIGN.CENTER)

# 项目英文名
add_text(slide1, Inches(0), Inches(4.2), W, Inches(0.5),
         "Taste of China — Interactive Food Culture Explorer", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# 团队信息
add_text(slide1, Inches(0), Inches(5.8), W, Inches(0.5),
         "团队成员：卢酉杰  王军帅", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

# 左右装饰小方块
add_rect(slide1, Inches(1.2), Inches(3.0), Inches(0.15), Inches(0.15), fill_color=RED)
add_rect(slide1, Inches(11.9), Inches(3.0), Inches(0.15), Inches(0.15), fill_color=RED)
add_rect(slide1, Inches(1.2), Inches(4.6), Inches(0.15), Inches(0.15), fill_color=GOLD)
add_rect(slide1, Inches(11.9), Inches(4.6), Inches(0.15), Inches(0.15), fill_color=GOLD)


# ============================================================
# 第2页：项目总述
# ============================================================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide2, BG_WARM)
add_top_bar(slide2)
add_bottom_bar(slide2)
add_page_number(slide2, 2)

add_decorated_title(slide2, "项目总述", "Project Overview")

# 项目简介卡片
intro_card = add_rounded_rect(slide2, Inches(0.6), Inches(1.5), Inches(6.0), Inches(2.2),
                               fill_color=WHITE, line_color=LIGHT_GRAY)
intro_card.line.width = Pt(1)
add_tag(slide2, Inches(0.9), Inches(1.65), "项目简介", fill_color=RED)
add_text(slide2, Inches(0.9), Inches(2.15), Inches(5.4), Inches(1.4),
         "以中国美食文化为主题的交互式探索网站，涵盖首页、地图探索、美食列表、详情页、挑战推荐五大模块，收录15道特色美食数据。用户可通过地图可视化、多维度搜索、智能推荐等多种方式发现和探索中国地域美食文化。",
         font_size=13, color=DARK)

# 创意来源卡片
origin_card = add_rounded_rect(slide2, Inches(6.9), Inches(1.5), Inches(6.0), Inches(2.2),
                                fill_color=WHITE, line_color=LIGHT_GRAY)
origin_card.line.width = Pt(1)
add_tag(slide2, Inches(7.2), Inches(1.65), "创意来源", fill_color=GOLD, text_color=DARK)
add_text(slide2, Inches(7.2), Inches(2.15), Inches(5.4), Inches(1.4),
         "中国拥有丰富多元的地域美食文化，但传统呈现方式缺乏交互性和趣味性。本项目将美食文化与Web技术结合，通过SVG地图可视化、智能匹配算法、成就挑战系统等创新交互，让用户在探索中感受中国美食的博大精深。",
         font_size=13, color=DARK)

# 五大模块卡片
modules = [
    ("🏠", "首页", "搜索 + 统计\n今日BOSS预览"),
    ("🗺️", "地图探索", "SVG中国地图\n省份交互 + 美食定位"),
    ("📋", "美食列表", "Grid布局\n二级联动筛选"),
    ("📖", "详情页", "数据驱动渲染\nURL参数解析"),
    ("🎯", "挑战推荐", "智能匹配算法\n成就记录系统"),
]

card_w = Inches(2.28)
card_h = Inches(2.6)
start_x = Inches(0.6)
gap = Inches(0.23)

for i, (icon, title, desc) in enumerate(modules):
    left = start_x + i * (card_w + gap)
    add_feature_card(slide2, left, Inches(4.1), card_w, card_h,
                     icon, title, desc, fill_color=WHITE, line_color=GOLD)

# 技术栈标签
add_text(slide2, Inches(0.6), Inches(6.9), Inches(12), Inches(0.35),
         "技术栈：HTML5 + CSS3 + JavaScript  |  数据存储：LocalStorage  |  地图渲染：SVG + GeoJSON  |  设计系统：CSS Variables + 响应式布局",
         font_size=11, color=GRAY, alignment=PP_ALIGN.LEFT)


# ============================================================
# 第3页：首页 — 搜索功能实现
# ============================================================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide3, BG_WARM)
add_top_bar(slide3)
add_bottom_bar(slide3)
add_page_number(slide3, 3)

add_decorated_title(slide3, "首页 · 搜索功能实现", "index.html — Real-time Multi-field Search")

# 左侧：功能说明
add_tag(slide3, Inches(0.6), Inches(1.4), "功能概述", fill_color=RED)

overview_items = [
    "▸ Hero区域：渐变背景 + 项目口号",
    "▸ 搜索模块：实时输入即搜索，0延迟",
    "▸ 统计网格：4列CSS Grid展示关键数据",
    "▸ 今日BOSS：基于日期哈希的确定性推荐",
    "▸ 快捷导航：一键跳转各功能模块",
]
add_bullet_points(slide3, Inches(0.6), Inches(1.85), Inches(5.5), Inches(2.5),
                  overview_items, font_size=13)

add_tag(slide3, Inches(0.6), Inches(4.0), "实现原理", fill_color=GOLD, text_color=DARK)

principle_items = [
    "▸ 监听input事件，触发实时过滤",
    "▸ 对foodBosses数组执行filter()",
    "▸ 6个维度模糊匹配：name / city / region / type / taste / desc",
    "▸ taste数组使用some()遍历匹配",
    "▸ 全部转小写后includes()比较，忽略大小写",
]
add_bullet_points(slide3, Inches(0.6), Inches(4.45), Inches(5.5), Inches(2.5),
                  principle_items, font_size=13)

# 右侧：核心代码
add_tag(slide3, Inches(6.5), Inches(1.4), "核心代码", fill_color=RED)

code_search = """const results = foodBosses.filter(food => {
    return food.name.toLowerCase().includes(keyword) ||
           food.city.toLowerCase().includes(keyword) ||
           food.region.toLowerCase().includes(keyword) ||
           food.type.toLowerCase().includes(keyword) ||
           food.taste.some(t =>
               t.toLowerCase().includes(keyword)) ||
           food.desc.toLowerCase().includes(keyword);
});"""
add_code_block(slide3, Inches(6.5), Inches(1.85), Inches(6.3), Inches(3.0), code_search)

# 搜索维度示意图
add_tag(slide3, Inches(6.5), Inches(5.1), "搜索维度", fill_color=GOLD, text_color=DARK)

dims = [("name", "名称"), ("city", "城市"), ("region", "地区"), ("type", "类型"), ("taste[]", "口味"), ("desc", "描述")]
for i, (key, label) in enumerate(dims):
    col = i % 3
    row = i // 3
    left = Inches(6.7) + col * Inches(2.1)
    top = Inches(5.55) + row * Inches(0.55)
    card = add_rounded_rect(slide3, left, top, Inches(1.8), Inches(0.42), fill_color=WHITE, line_color=RED)
    card.line.width = Pt(1.2)
    add_text(slide3, left + Inches(0.05), top + Inches(0.02), Inches(1.7), Inches(0.38),
             f"{key}  →  {label}", font_size=11, color=DARK, alignment=PP_ALIGN.CENTER)


# ============================================================
# 第4页：地图页 — SVG渲染原理
# ============================================================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide4, BG_WARM)
add_top_bar(slide4)
add_bottom_bar(slide4)
add_page_number(slide4, 4)

add_decorated_title(slide4, "地图页 · SVG地图渲染原理", "map.html — GeoJSON to SVG Path Pipeline")

# 技术流程图
flow_steps = ["GeoJSON\n原始数据", "坐标投影\nproject()", "SVG Path\nd属性拼接", "DOM渲染\n交互绑定"]
for i, step in enumerate(flow_steps):
    left = Inches(0.8) + i * Inches(3.1)
    card = add_rounded_rect(slide4, left, Inches(1.5), Inches(2.4), Inches(1.0),
                             fill_color=WHITE, line_color=GOLD)
    card.line.width = Pt(1.5)
    add_text(slide4, left + Inches(0.1), Inches(1.55), Inches(2.2), Inches(0.9),
             step, font_size=13, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 3:
        add_text(slide4, left + Inches(2.4), Inches(1.7), Inches(0.6), Inches(0.5),
                 "→", font_size=24, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# 坐标投影算法代码
add_tag(slide4, Inches(0.6), Inches(2.8), "坐标投影算法", fill_color=RED)

code_project = """function project(lng, lat) {
    const x = ((lng - MAP_LNG_MIN) /
        (MAP_LNG_MAX - MAP_LNG_MIN)) * MAP_W;
    const y = ((MAP_LAT_MAX - lat) /
        (MAP_LAT_MAX - MAP_LAT_MIN)) * MAP_H;
    return [x.toFixed(2), y.toFixed(2)];
}"""
add_code_block(slide4, Inches(0.6), Inches(3.25), Inches(5.8), Inches(2.7), code_project)

# GeoJSON处理说明
add_tag(slide4, Inches(6.8), Inches(2.8), "GeoJSON处理流程", fill_color=GOLD, text_color=DARK)

geo_items = [
    "① 加载china_geo_data.js中的GeoJSON对象",
    "② 遍历features数组获取各省级行政区",
    "③ 对每个feature的coordinates递归处理",
    "④ MultiPolygon类型逐层遍历坐标数组",
    "⑤ 每组坐标经project()转换得[x, y]",
    "⑥ 拼接 M x y L x y ... Z 生成SVG path",
    "⑦ 创建<path>元素设置d属性与样式",
    "⑧ 绑定hover/click事件实现交互",
]
add_bullet_points(slide4, Inches(6.8), Inches(3.25), Inches(6.0), Inches(3.5),
                  geo_items, font_size=12, spacing=Pt(4))

# 底部运行效果说明
add_tag(slide4, Inches(0.6), Inches(6.2), "运行效果", fill_color=RED)
add_text(slide4, Inches(0.6), Inches(6.6), Inches(12), Inches(0.4),
         "完整中国SVG地图  ·  省份hover高亮 + Tooltip信息  ·  安徽省特殊脉冲动画 + 渐变填充  ·  美食标记点定位",
         font_size=12, color=DARK)


# ============================================================
# 第5页：地图页 — 搜索与视觉效果
# ============================================================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide5, BG_WARM)
add_top_bar(slide5)
add_bottom_bar(slide5)
add_page_number(slide5, 5)

add_decorated_title(slide5, "地图页 · 搜索与视觉效果", "map.html — Search & Visual Effects")

# 地图搜索说明
add_tag(slide5, Inches(0.6), Inches(1.4), "地图搜索", fill_color=RED)

search_items = [
    "▸ 输入关键词实时过滤省份",
    "▸ 匹配省份保持正常显示 + 高亮边框",
    "▸ 非匹配省份透明度降低至0.3",
    "▸ 搜索结果数量实时显示",
]
add_bullet_points(slide5, Inches(0.6), Inches(1.85), Inches(5.5), Inches(1.8),
                  search_items, font_size=13)

# 搜索代码
code_map_search = """// 地图搜索过滤逻辑
provinces.forEach(p => {
    const match = p.name.includes(keyword) ||
        p.foods.some(f => f.includes(keyword));
    p.element.style.opacity = match ? '1' : '0.3';
    p.element.style.stroke = match ?
        '#C41E24' : '#ccc';
});"""
add_code_block(slide5, Inches(0.6), Inches(3.3), Inches(5.5), Inches(2.2), code_map_search)

# 视觉效果详解
add_tag(slide5, Inches(6.5), Inches(1.4), "视觉效果", fill_color=GOLD, text_color=DARK)

effects = [
    ("省份Hover", "transform: scale(1.02)\nbox-shadow阴影增强\nstroke颜色变红", RED),
    ("安徽脉冲", "CSS @keyframes pulse\n渐变填充 animate\n0%→50%→100% 透明度", GOLD),
    ("毛玻璃Tooltip", "backdrop-filter: blur(10px)\nrgba半透明背景\n跟随鼠标移动", RED_DARK),
    ("美食标记点", "定位到省份中心\n弹跳动画效果\n点击跳转详情页", DARK_LIGHT),
]

for i, (title, desc, color) in enumerate(effects):
    col = i % 2
    row = i // 2
    left = Inches(6.5) + col * Inches(3.3)
    top = Inches(1.85) + row * Inches(2.1)
    card = add_rounded_rect(slide5, left, top, Inches(3.0), Inches(1.9),
                             fill_color=WHITE, line_color=color)
    card.line.width = Pt(2)
    add_text(slide5, left + Inches(0.15), top + Inches(0.1), Inches(2.7), Inches(0.35),
             title, font_size=14, color=color, bold=True)
    add_text(slide5, left + Inches(0.15), top + Inches(0.5), Inches(2.7), Inches(1.3),
             desc, font_size=11, color=GRAY)

# 底部CSS代码片段
code_css_effect = """/* 安徽脉冲动画 */
@keyframes pulse {
    0%   { opacity: 0.7; transform: scale(1); }
    50%  { opacity: 1;   transform: scale(1.02); }
    100% { opacity: 0.7; transform: scale(1); }
}"""
add_code_block(slide5, Inches(0.6), Inches(5.7), Inches(5.5), Inches(1.3), code_css_effect)


# ============================================================
# 第6页：列表页 — 布局原理
# ============================================================
slide6 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide6, BG_WARM)
add_top_bar(slide6)
add_bottom_bar(slide6)
add_page_number(slide6, 6)

add_decorated_title(slide6, "列表页 · 布局原理", "list.html — CSS Grid + Sticky + Frosted Glass")

# CSS Grid响应式
add_tag(slide6, Inches(0.6), Inches(1.4), "CSS Grid 响应式布局", fill_color=RED)

grid_items = [
    "▸ >1024px：3列网格，信息密度最优",
    "▸ 768px~1024px：2列网格，平板适配",
    "▸ <768px：1列布局，移动端友好",
    "▸ grid-template-columns + auto-fill + minmax()",
]
add_bullet_points(slide6, Inches(0.6), Inches(1.85), Inches(5.5), Inches(1.6),
                  grid_items, font_size=13)

# 粘性面板代码
add_tag(slide6, Inches(0.6), Inches(3.3), "粘性控制面板 + 毛玻璃", fill_color=GOLD, text_color=DARK)

code_sticky = """.control-panel-sticky {
    position: sticky;
    top: 57px;
    z-index: 150;
    background: rgba(255, 255, 255, 0.85);
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border-bottom: 1px solid rgba(212,168,67,0.2);
}"""
add_code_block(slide6, Inches(0.6), Inches(3.75), Inches(6.0), Inches(2.5), code_sticky)

# 右侧：其他特性
add_tag(slide6, Inches(7.0), Inches(1.4), "交互特性", fill_color=RED)

features_right = [
    ("二级联动筛选", "地区选择 → 动态更新省份列表\n省级联动，选择地区后\n自动过滤该地区下的省份"),
    ("实时搜索", "input事件监听 + filter()\n实时显示匹配结果数量badge\n支持名称/城市/口味等多字段"),
    ("毛玻璃原理", "backdrop-filter: blur(20px)\n对元素背后区域进行模糊\n配合半透明背景实现通透感"),
]

for i, (title, desc) in enumerate(features_right):
    top = Inches(1.85) + i * Inches(1.7)
    card = add_rounded_rect(slide6, Inches(7.0), top, Inches(5.8), Inches(1.5),
                             fill_color=WHITE, line_color=GOLD)
    card.line.width = Pt(1.2)
    add_text(slide6, Inches(7.2), top + Inches(0.1), Inches(5.4), Inches(0.3),
             title, font_size=14, color=RED, bold=True)
    add_text(slide6, Inches(7.2), top + Inches(0.45), Inches(5.4), Inches(0.9),
             desc, font_size=12, color=GRAY)

# 底部说明
add_text(slide6, Inches(7.0), Inches(6.5), Inches(5.8), Inches(0.4),
         "关键：sticky定位让控制面板在滚动时始终可见，提升操作效率",
         font_size=11, color=DARK)


# ============================================================
# 第7页：详情页 — 动态渲染
# ============================================================
slide7 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide7, BG_WARM)
add_top_bar(slide7)
add_bottom_bar(slide7)
add_page_number(slide7, 7)

add_decorated_title(slide7, "详情页 · 数据驱动动态渲染", "detail.html — URL Params + Dynamic Rendering")

# 页面跳转流程
add_tag(slide7, Inches(0.6), Inches(1.4), "页面跳转流程", fill_color=RED)

flow_items = ["列表页\nlist.html", "点击卡片\n?a=id", "详情页\ndetail.html"]
for i, item in enumerate(flow_items):
    left = Inches(0.8) + i * Inches(3.5)
    card = add_rounded_rect(slide7, left, Inches(1.85), Inches(2.5), Inches(0.8),
                             fill_color=WHITE, line_color=GOLD)
    card.line.width = Pt(1.5)
    add_text(slide7, left + Inches(0.1), Inches(1.9), Inches(2.3), Inches(0.7),
             item, font_size=13, color=DARK, bold=True, alignment=PP_ALIGN.CENTER)
    if i < 2:
        add_text(slide7, left + Inches(2.5), Inches(1.95), Inches(0.8), Inches(0.5),
                 "→", font_size=22, color=RED, bold=True, alignment=PP_ALIGN.CENTER)

# URL参数解析代码
add_tag(slide7, Inches(0.6), Inches(3.0), "URL参数解析", fill_color=RED)

code_url = """// 从URL获取美食ID
const urlParams = new URLSearchParams(
    window.location.search
);
const foodId = parseInt(urlParams.get('id'));

// 根据ID查找数据
const food = foodBosses.find(
    f => f.id === foodId
);"""
add_code_block(slide7, Inches(0.6), Inches(3.45), Inches(5.5), Inches(2.3), code_url)

# 右侧：渲染特性
add_tag(slide7, Inches(6.5), Inches(3.0), "渲染特性", fill_color=GOLD, text_color=DARK)

render_features = [
    "▸ URLSearchParams解析查询字符串",
    "▸ find()方法根据id精确查找数据对象",
    "▸ 数据驱动：从data.js读取后动态生成DOM",
    "▸ 渐变色头部：使用美食配色方案",
    "▸ 两栏信息网格：CSS Grid 2列布局",
    "▸ 标签展示：口味、地区、类型等标签",
    "▸ 响应式适配：移动端切换为单列",
    "▸ 返回按钮：history.back()返回列表",
]
add_bullet_points(slide7, Inches(6.5), Inches(3.45), Inches(6.0), Inches(3.0),
                  render_features, font_size=13, spacing=Pt(5))

# 底部关键点
add_rounded_rect(slide7, Inches(0.6), Inches(6.0), Inches(12.0), Inches(0.7),
                  fill_color=RGBColor(0xFF, 0xF0, 0xE0), line_color=GOLD)
add_text(slide7, Inches(0.9), Inches(6.1), Inches(11.4), Inches(0.5),
         "💡 核心思路：列表页通过URL传递参数 → 详情页解析参数 → 查找数据 → 动态渲染页面内容，实现页面间的数据流转",
         font_size=13, color=DARK, bold=True)


# ============================================================
# 第8页：挑战页 — 智能推荐算法
# ============================================================
slide8 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide8, BG_WARM)
add_top_bar(slide8)
add_bottom_bar(slide8)
add_page_number(slide8, 8)

add_decorated_title(slide8, "挑战页 · 智能推荐算法", "challenge.html — Multi-dimensional Weighted Scoring")

# 算法代码
add_tag(slide8, Inches(0.6), Inches(1.4), "核心算法", fill_color=RED)

code_algo = """function calculateMatch(food, conditions) {
    let score = 60;  // 基础分
    // 口味匹配 +15
    if (food.taste.some(t =>
        conditions.taste.includes(t)))
        score += 15;
    // 场景匹配 +10
    if (food.scene === conditions.scene)
        score += 10;
    // 心情匹配 +10
    if (food.mood === conditions.mood)
        score += 10;
    // 天气匹配 +5
    if (food.weather === conditions.weather)
        score += 5;
    // 安徽加分 +5
    if (food.province === '安徽')
        score += 5;
    score = Math.min(score, 100);
    return { score, reasons };
}"""
add_code_block(slide8, Inches(0.6), Inches(1.85), Inches(6.0), Inches(4.5), code_algo)

# 右侧：评分维度可视化
add_tag(slide8, Inches(7.0), Inches(1.4), "评分维度权重", fill_color=GOLD, text_color=DARK)

dimensions = [
    ("基础分", 60, "每道美食的起始分数", RED_LIGHT),
    ("口味匹配", 15, "用户口味偏好与美食口味交集", RED),
    ("场景匹配", 10, "用餐场景是否吻合", GOLD),
    ("心情匹配", 10, "当前心情与美食适配度", GOLD),
    ("天气匹配", 5, "天气对美食选择的影响", DARK_LIGHT),
    ("安徽加分", 5, "本地美食额外加分", DARK_LIGHT),
]

for i, (name, score, desc, color) in enumerate(dimensions):
    top = Inches(1.85) + i * Inches(0.75)
    # 标签
    add_text(slide8, Inches(7.0), top, Inches(1.5), Inches(0.35),
             name, font_size=12, color=DARK, bold=True)
    # 进度条背景
    add_rounded_rect(slide8, Inches(8.5), top + Inches(0.05), Inches(3.5), Inches(0.25),
                      fill_color=RGBColor(0xE8, 0xE8, 0xE8))
    # 进度条填充
    bar_width = 3.5 * (score / 60)
    add_rounded_rect(slide8, Inches(8.5), top + Inches(0.05), Inches(bar_width), Inches(0.25),
                      fill_color=color)
    # 分数
    add_text(slide8, Inches(12.1), top, Inches(0.8), Inches(0.35),
             f"+{score}", font_size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)
    # 说明
    add_text(slide8, Inches(7.0), top + Inches(0.32), Inches(5.5), Inches(0.3),
             desc, font_size=10, color=GRAY)

# 今日BOSS
add_tag(slide8, Inches(7.0), Inches(6.3), "今日BOSS", fill_color=RED)
add_text(slide8, Inches(8.2), Inches(6.3), Inches(4.5), Inches(0.35),
         "基于日期字符串的确定性哈希，保证每日推荐一致", font_size=11, color=DARK)


# ============================================================
# 第9页：挑战记录与成就系统
# ============================================================
slide9 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide9, BG_WARM)
add_top_bar(slide9)
add_bottom_bar(slide9)
add_page_number(slide9, 9)

add_decorated_title(slide9, "挑战记录与成就系统", "Challenge Records & Achievement System")

# LocalStorage持久化
add_tag(slide9, Inches(0.6), Inches(1.4), "LocalStorage 持久化", fill_color=RED)

code_storage = """// 保存挑战记录
function saveChallengeRecord(record) {
    const records = getChallengeRecords();
    records.push({
        foodId: record.foodId,
        score: record.score,
        conditions: record.conditions,
        timestamp: Date.now()
    });
    localStorage.setItem(
        'bossChallengeRecords',
        JSON.stringify(records)
    );
}"""
add_code_block(slide9, Inches(0.6), Inches(1.85), Inches(6.0), Inches(2.8), code_storage)

# 数据统计
add_tag(slide9, Inches(6.8), Inches(1.4), "数据统计函数", fill_color=GOLD, text_color=DARK)

stats_items = [
    "▸ getChallengeRecords()：获取全部记录",
    "▸ getChallengeStats()：统计总次数/最爱口味/最常地区等",
    "▸ getAchievements()：遍历7种成就判定条件",
    "▸ 数据结构：{ foodId, score, conditions, timestamp }",
    "▸ 重试机制：retryRecommendation()循环推荐列表",
]
add_bullet_points(slide9, Inches(6.8), Inches(1.85), Inches(6.0), Inches(2.5),
                  stats_items, font_size=12, spacing=Pt(6))

# 成就系统
add_tag(slide9, Inches(0.6), Inches(4.9), "成就徽章系统", fill_color=RED)

achievements = [
    ("🏆", "初出茅庐", "完成首次挑战"),
    ("🍜", "口味探索者", "尝试3种不同口味"),
    ("🗺️", "区域达人", "探索3个不同地区"),
    ("⭐", "满分BOSS", "获得100分匹配"),
    ("🔥", "连续挑战", "连续挑战5次"),
    ("🎯", "精准推荐", "匹配度≥90分"),
    ("👑", "美食王者", "解锁全部成就"),
]

for i, (icon, name, condition) in enumerate(achievements):
    col = i % 4
    row = i // 4
    left = Inches(0.6) + col * Inches(3.15)
    top = Inches(5.4) + row * Inches(0.9)
    card = add_rounded_rect(slide9, left, top, Inches(2.9), Inches(0.75),
                             fill_color=WHITE, line_color=GOLD)
    card.line.width = Pt(1.2)
    add_text(slide9, left + Inches(0.1), top + Inches(0.05), Inches(0.4), Inches(0.35),
             icon, font_size=16, alignment=PP_ALIGN.CENTER)
    add_text(slide9, left + Inches(0.5), top + Inches(0.05), Inches(1.2), Inches(0.3),
             name, font_size=12, color=DARK, bold=True)
    add_text(slide9, left + Inches(0.5), top + Inches(0.38), Inches(2.2), Inches(0.3),
             condition, font_size=10, color=GRAY)


# ============================================================
# 第10页：结尾页
# ============================================================
slide10 = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide10, DARK)

# 装饰
add_rect(slide10, Inches(0), Inches(0), W, Inches(0.12), fill_color=RED)
add_rect(slide10, Inches(0), Inches(7.38), W, Inches(0.12), fill_color=RED)

# 金色装饰线
add_rect(slide10, Inches(3.5), Inches(2.2), Inches(6.333), Inches(0.03), fill_color=GOLD)
add_rect(slide10, Inches(3.5), Inches(5.2), Inches(6.333), Inches(0.03), fill_color=GOLD)

# 谢谢观看
add_text(slide10, Inches(0), Inches(2.5), W, Inches(1.0),
         "谢谢观看", font_size=54, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# 项目名
add_text(slide10, Inches(0), Inches(3.6), W, Inches(0.6),
         "舌尖BOSS — 中国美食文化探索网站", font_size=22, color=GOLD, alignment=PP_ALIGN.CENTER)

# 团队
add_text(slide10, Inches(0), Inches(4.5), W, Inches(0.5),
         "团队成员：卢酉杰  王军帅", font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)

# 装饰方块
add_rect(slide10, Inches(1.2), Inches(3.0), Inches(0.15), Inches(0.15), fill_color=GOLD)
add_rect(slide10, Inches(11.9), Inches(3.0), Inches(0.15), Inches(0.15), fill_color=GOLD)
add_rect(slide10, Inches(1.2), Inches(4.0), Inches(0.15), Inches(0.15), fill_color=RED)
add_rect(slide10, Inches(11.9), Inches(4.0), Inches(0.15), Inches(0.15), fill_color=RED)

# 底部标语
add_text(slide10, Inches(0), Inches(5.5), W, Inches(0.4),
         "用代码品味中国  ·  用交互探索美食", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# ========== 保存 ==========
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "舌尖BOSS项目汇报.pptx")
prs.save(output_path)
print(f"PPTX saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
